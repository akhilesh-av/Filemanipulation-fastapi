from fastapi import FastAPI, File, UploadFile, HTTPException, Form
from fastapi.responses import FileResponse, StreamingResponse
from PyPDF2 import PdfMerger, PdfReader, PdfWriter
from docx import Document as DocxDocument
from PIL import Image
import io
import zipfile
import comtypes.client
import os
from typing import Optional
from fastapi import FastAPI, File, UploadFile, Form, HTTPException, Depends
from fastapi.responses import FileResponse, StreamingResponse
from PyPDF2 import PdfMerger, PdfReader, PdfWriter
from docx import Document as DocxDocument
from PIL import Image
import io
import zipfile
import os
import pdfplumber
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter, landscape, A4
from reportlab.lib.utils import ImageReader
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
import fitz  # PyMuPDF
import openpyxl
from openpyxl.drawing.image import Image as OpenpyxlImage
from typing import List, Optional
from reportlab.lib.colors import black, lightgrey
from reportlab.lib.units import inch
from pdf2docx import Converter


app = FastAPI()


@app.post("/merge_pdfs/")
async def merge_pdfs(file_1: UploadFile = File(...), file_2: UploadFile = File(...)):
    try:
        # Define paths in the root directory
        root_dir = os.path.dirname(os.path.abspath(__file__))
        file_1_path = os.path.join(root_dir, file_1.filename)
        file_2_path = os.path.join(root_dir, file_2.filename)
        output_filename = "merged_output.pdf"
        output_path = os.path.join(root_dir, output_filename)

        # Save the uploaded files to the root directory
        with open(file_1_path, "wb") as f:
            f.write(await file_1.read())
        with open(file_2_path, "wb") as f:
            f.write(await file_2.read())

        merger = PdfMerger()
        merger.append(file_1_path)
        merger.append(file_2_path)
        merger.write(output_path)
        merger.close()

        return FileResponse(output_path, filename=output_filename, media_type='application/pdf')

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

    finally:
        # Clean up the input files, but keep the output file
        if os.path.exists(file_1_path):
            os.remove(file_1_path)
        if os.path.exists(file_2_path):
            os.remove(file_2_path)



@app.post("/merge_docx/")
async def merge_docx(file_1: UploadFile = File(...), file_2: UploadFile = File(...)):
    try:
        # Define paths in the root directory
        root_dir = os.path.dirname(os.path.abspath(__file__))
        file_1_path = os.path.join(root_dir, file_1.filename)
        file_2_path = os.path.join(root_dir, file_2.filename)
        output_filename = "merged_output.docx"
        output_path = os.path.join(root_dir, output_filename)

        # Save the uploaded files to the root directory
        with open(file_1_path, "wb") as f:
            f.write(await file_1.read())
        with open(file_2_path, "wb") as f:
            f.write(await file_2.read())

        merged_document = DocxDocument(file_1_path)
        doc2 = DocxDocument(file_2_path)

        for element in doc2.element.body:
            merged_document.element.body.append(element)

        merged_document.save(output_path)

        return FileResponse(output_path, filename=output_filename, media_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document')

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

    finally:
        # Clean up the input files, but keep the output file
        if os.path.exists(file_1_path):
            os.remove(file_1_path)
        if os.path.exists(file_2_path):
            os.remove(file_2_path)



@app.post("/merge_img/")
async def merge_img(file_1: UploadFile = File(...), file_2: UploadFile = File(...), side_by_side: Optional[bool] = Form(True)):
    try:
        # Define paths in the root directory
        root_dir = os.path.dirname(os.path.abspath(__file__))
        file_1_path = os.path.join(root_dir, file_1.filename)
        file_2_path = os.path.join(root_dir, file_2.filename)
        output_filename = "merged_output.png"
        output_path = os.path.join(root_dir, output_filename)

        # Save the uploaded files to the root directory
        with open(file_1_path, "wb") as f:
            f.write(await file_1.read())
        with open(file_2_path, "wb") as f:
            f.write(await file_2.read())

        img1 = Image.open(file_1_path)
        img2 = Image.open(file_2_path)

        if side_by_side:
            merged_width = img1.width + img2.width
            merged_height = max(img1.height, img2.height)
            merged_image = Image.new("RGBA", (merged_width, merged_height))
            merged_image.paste(img1, (0, 0))
            merged_image.paste(img2, (img1.width, 0))
        else:
            merged_width = max(img1.width, img2.width)
            merged_height = img1.height + img2.height
            merged_image = Image.new("RGBA", (merged_width, merged_height))
            merged_image.paste(img1, (0, 0))
            merged_image.paste(img2, (0, img1.height))

        merged_image.save(output_path)

        return FileResponse(output_path, filename=output_filename, media_type='image/png')

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

    finally:
        # Clean up the input files, but keep the output file
        if os.path.exists(file_1_path):
            os.remove(file_1_path)
        if os.path.exists(file_2_path):
            os.remove(file_2_path)




@app.post("/split_pdf/")
async def split_pdf(file: UploadFile = File(...)):
    try:
        # Read the uploaded file into a BytesIO object
        file_content = await file.read()
        pdf_stream = io.BytesIO(file_content)

        # Create a PdfReader object
        pdf_reader = PdfReader(pdf_stream)

        zip_stream = io.BytesIO()
        with zipfile.ZipFile(zip_stream, 'w') as zip_file:
            for page_number in range(len(pdf_reader.pages)):
                pdf_writer = PdfWriter()
                pdf_writer.add_page(pdf_reader.pages[page_number])

                pdf_page_stream = io.BytesIO()
                pdf_writer.write(pdf_page_stream)
                pdf_page_stream.seek(0)

                output_pdf_filename = f"page_{page_number + 1}.pdf"
                zip_file.writestr(output_pdf_filename, pdf_page_stream.getvalue())

        zip_stream.seek(0)

        return StreamingResponse(zip_stream, media_type="application/zip", headers={"Content-Disposition": "attachment; filename=split_pages.zip"})

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))





@app.post("/trim_pdf/")
async def trim_pdf(file: UploadFile = File(...), start_page: int = Form(...), end_page: int = Form(...)):
    try:
        # Read the uploaded file into a BytesIO object
        file_content = await file.read()
        pdf_stream = io.BytesIO(file_content)

        # Create a PdfReader object
        pdf_reader = PdfReader(pdf_stream)
        total_pages = len(pdf_reader.pages)

        if start_page < 1 or end_page > total_pages or start_page > end_page:
            raise HTTPException(status_code=400, detail=f"Please provide valid page numbers between 1 and {total_pages}.")

        pdf_writer = PdfWriter()

        for page_number in range(start_page - 1, end_page):
            pdf_writer.add_page(pdf_reader.pages[page_number])

        trimmed_pdf_stream = io.BytesIO()
        pdf_writer.write(trimmed_pdf_stream)
        trimmed_pdf_stream.seek(0)

        return StreamingResponse(trimmed_pdf_stream, media_type='application/pdf', headers={"Content-Disposition": "attachment; filename=trimmed_output.pdf"})

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))



@app.post("/flatten_pdf/")
async def flatten_pdf(file: UploadFile = File(...)):
    try:
        # Read the uploaded file into a BytesIO object
        file_content = await file.read()
        pdf_stream = io.BytesIO(file_content)

        # Create a PdfReader object
        pdf_reader = PdfReader(pdf_stream)
        pdf_writer = PdfWriter()

        for page in pdf_reader.pages:
            pdf_writer.add_page(page)
            if '/Annots' in page:
                page['/Annots'] = []

        flattened_pdf_stream = io.BytesIO()
        pdf_writer.write(flattened_pdf_stream)
        flattened_pdf_stream.seek(0)

        return StreamingResponse(flattened_pdf_stream, media_type='application/pdf', headers={"Content-Disposition": "attachment; filename=flattened_output.pdf"})

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))



@app.post("/add_watermark/")
async def add_watermark(file: UploadFile = File(...), watermark_text: str = Form(...)):
    try:
        # Read the uploaded file into a BytesIO object
        file_content = await file.read()
        pdf_stream = io.BytesIO(file_content)

        # Create a PdfReader object
        pdf_reader = PdfReader(pdf_stream)
        pdf_writer = PdfWriter()

        packet = io.BytesIO()
        can = canvas.Canvas(packet, pagesize=letter)

        for _ in range(len(pdf_reader.pages)):
            can.saveState()
            can.setFont("Helvetica", 40)
            can.setFillColorRGB(0, 0, 0, 0.1)
            can.drawString(200, 300, watermark_text)
            can.restoreState()
            can.showPage()

        can.save()
        packet.seek(0)

        watermark_pdf = PdfReader(packet)

        for i in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[i]
            watermark_page = watermark_pdf.pages[0]
            page.merge_page(watermark_page)
            pdf_writer.add_page(page)

        output_stream = io.BytesIO()
        pdf_writer.write(output_stream)
        output_stream.seek(0)

        return StreamingResponse(output_stream, media_type='application/pdf', headers={"Content-Disposition": "attachment; filename=watermarked_output.pdf"})

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))



@app.post("/add_page_numbers/")
async def add_page_numbers(file: UploadFile = File(...)):
    try:
        # Read the uploaded file into a BytesIO object
        file_content = await file.read()
        pdf_stream = io.BytesIO(file_content)

        # Create a PdfReader object
        pdf_reader = PdfReader(pdf_stream)
        pdf_writer = PdfWriter()

        packet = io.BytesIO()
        can = canvas.Canvas(packet, pagesize=letter)

        for page_num in range(len(pdf_reader.pages)):
            can.drawString(270, 10, str(page_num + 1))
            can.showPage()

        can.save()
        packet.seek(0)

        page_numbers_pdf = PdfReader(packet)

        for i in range(len(pdf_reader.pages)):
            original_page = pdf_reader.pages[i]
            number_page = page_numbers_pdf.pages[i]
            original_page.merge_page(number_page)
            pdf_writer.add_page(original_page)

        output_stream = io.BytesIO()
        pdf_writer.write(output_stream)
        output_stream.seek(0)

        return StreamingResponse(output_stream, media_type='application/pdf', headers={"Content-Disposition": "attachment; filename=numbered_output.pdf"})

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/rotate_pdf/")
async def rotate_pdf(file: UploadFile = File(...), rotation_angle: int = Form(...)):
    try:
        if rotation_angle not in [90, 180, 270]:
            raise HTTPException(status_code=400, detail="Invalid rotation angle. Please use 90, 180, or 270 degrees.")

        # Read the uploaded file into a BytesIO object
        file_content = await file.read()
        pdf_stream = io.BytesIO(file_content)

        pdf_reader = PdfReader(pdf_stream)
        pdf_writer = PdfWriter()

        for page in pdf_reader.pages:
            page.rotate(rotation_angle)
            pdf_writer.add_page(page)

        output_stream = io.BytesIO()
        pdf_writer.write(output_stream)
        output_stream.seek(0)

        return StreamingResponse(output_stream, media_type='application/pdf', headers={"Content-Disposition": "attachment; filename=rotated_output.pdf"})

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/delete_pdf_page/")
async def delete_pdf_page(file: UploadFile = File(...), page_to_delete: int = Form(...)):
    try:
        # Read the uploaded file into a BytesIO object
        file_content = await file.read()
        pdf_stream = io.BytesIO(file_content)

        pdf_reader = PdfReader(pdf_stream)
        num_pages = len(pdf_reader.pages)

        if page_to_delete < 1 or page_to_delete > num_pages:
            raise HTTPException(status_code=400, detail=f"Invalid page number. Please enter a number between 1 and {num_pages}.")

        page_index_to_delete = page_to_delete - 1
        pdf_writer = PdfWriter()

        for i in range(num_pages):
            if i != page_index_to_delete:
                pdf_writer.add_page(pdf_reader.pages[i])

        output_stream = io.BytesIO()
        pdf_writer.write(output_stream)
        output_stream.seek(0)

        return StreamingResponse(output_stream, media_type='application/pdf', headers={"Content-Disposition": "attachment; filename=updated_output.pdf"})

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    
    
@app.post("/pdf_to_excel/")
async def pdf_to_excel(file: UploadFile = File(...)):
    try:
        all_text = []
        with pdfplumber.open(await file.read()) as pdf:
            for page in pdf.pages:
                tables = page.extract_tables()
                for table in tables:
                    if table:
                        df = pd.DataFrame(table[1:], columns=table[0])
                        all_text.append(df)

        if not all_text:
            raise HTTPException(status_code=204, detail="No tables found in the PDF.")

        result_df = pd.concat(all_text, ignore_index=True)

        excel_buffer = io.BytesIO()
        with pd.ExcelWriter(excel_buffer, engine='openpyxl') as writer:
            result_df.to_excel(writer, index=False, sheet_name='Sheet1')

        excel_buffer.seek(0)

        return StreamingResponse(excel_buffer, media_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', headers={"Content-Disposition": "attachment; filename=converted_output.xlsx"})

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/excel_to_pdf/")
async def excel_to_pdf(file: UploadFile = File(...)):
    try:
        excel_data = pd.read_excel(await file.read())

        pdf_buffer = io.BytesIO()
        pdf = canvas.Canvas(pdf_buffer, pagesize=landscape(A4))
        width, height = landscape(A4)

        x_offset, y_offset = 0.5 * inch, height - inch
        row_height = 20
        col_width = width / (len(excel_data.columns) + 1)

        pdf.setFillColor(lightgrey)
        pdf.rect(x_offset, y_offset, width - (2 * inch), row_height, fill=True, stroke=False)
        pdf.setFillColor(black)

        pdf.setFont("Helvetica-Bold", 10)
        for i, col_name in enumerate(excel_data.columns):
            pdf.drawString(x_offset + i * col_width, y_offset + 5, str(col_name))

        y_offset -= row_height

        pdf.setFont("Helvetica", 10)
        for _, row in excel_data.iterrows():
            for i, cell in enumerate(row):
                pdf.drawString(x_offset + i * col_width, y_offset + 5, str(cell))
            y_offset -= row_height

            if y_offset < inch:
                pdf.showPage()
                y_offset = height - inch

        pdf.save()
        pdf_buffer.seek(0)

        return StreamingResponse(pdf_buffer, media_type='application/pdf', headers={"Content-Disposition": "attachment; filename=excel_to_pdf_output.pdf"})

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/pdf_to_pptx/")
async def pdf_to_pptx(file: UploadFile = File(...)):
    try:
        pptx_buffer = io.BytesIO()
        prs = Presentation()

        pdf_document = fitz.open(stream=await file.read(), filetype="pdf")

        for i in range(len(pdf_document)):
            page = pdf_document[i]
            pix = page.get_pixmap()

            img_buffer = io.BytesIO()
            img_bytes = pix.tobytes("png")
            img_buffer.write(img_bytes)
            img_buffer.seek(0)

            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.add_picture(img_buffer, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

        prs.save(pptx_buffer)
        pptx_buffer.seek(0)

        return StreamingResponse(pptx_buffer, media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation', headers={"Content-Disposition": "attachment; filename=converted_presentation.pptx"})

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/word_to_pdf/")
async def word_to_pdf(file: UploadFile = File(...)):
    try:
        word_path = f"temp_{file.filename}"
        with open(word_path, 'wb') as temp_word_file:
            temp_word_file.write(await file.read())

        pdf_path = f"{os.path.splitext(word_path)[0]}.pdf"

        word_app = comtypes.client.CreateObject('Word.Application')
        word_app.Visible = False
        doc = word_app.Documents.Open(os.path.abspath(word_path))
        doc.SaveAs(os.path.abspath(pdf_path), FileFormat=17)
        doc.Close()
        word_app.Quit()

        with open(pdf_path, 'rb') as pdf_file:
            pdf_data = pdf_file.read()

        os.remove(word_path)
        os.remove(pdf_path)

        return StreamingResponse(io.BytesIO(pdf_data), media_type='application/pdf', headers={"Content-Disposition": f'attachment; filename="{os.path.basename(pdf_path)}"'})

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/pdf_to_word/")
async def pdf_to_word(file: UploadFile = File(...)):
    try:
        pdf_path = f"temp_{file.filename}"
        with open(pdf_path, 'wb') as temp_pdf_file:
            temp_pdf_file.write(await file.read())

        word_path = f"{os.path.splitext(pdf_path)[0]}.docx"

        cv = Converter(pdf_path)
        cv.convert(word_path, start=0, end=None)
        cv.close()

        with open(word_path, 'rb') as word_file:
            word_data = word_file.read()

        os.remove(pdf_path)
        os.remove(word_path)

        return StreamingResponse(io.BytesIO(word_data), media_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document', headers={"Content-Disposition": f'attachment; filename="{os.path.basename(word_path)}"'})

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/compress_image/")
async def compress_image(file: UploadFile = File(...)):
    try:
        image = Image.open(await file.read())

        if image.mode in ("RGBA", "P"):
            image = image.convert("RGB")

        compressed_image_io = io.BytesIO()
        image.save(compressed_image_io, 'JPEG', optimize=True, quality=10)
        compressed_image_io.seek(0)

        return StreamingResponse(compressed_image_io, media_type='image/jpeg', headers={"Content-Disposition": "attachment; filename=compressed_image.jpg"})

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/compress_pptx_images/")
async def compress_pptx_images(file: UploadFile = File(...)):
    try:
        prs = Presentation(await file.read())
        temp_image_folder = "temp_images"
        os.makedirs(temp_image_folder, exist_ok=True)

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:
                    original_image_path = os.path.join(temp_image_folder, "original_image.jpg")
                    with open(original_image_path, "wb") as f:
                        f.write(shape.image.blob)

                    compressed_image_path = os.path.join(temp_image_folder, "compressed_image.jpg")
                    with Image.open(original_image_path) as img:
                        img.save(compressed_image_path, optimize=True, quality=50)

                    left, top, width, height = shape.left, shape.top, shape.width, shape.height
                    slide.shapes._spTree.remove(shape._element)
                    slide.shapes.add_picture(compressed_image_path, left, top, width, height)

        pptx_io = io.BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)

        for filename in os.listdir(temp_image_folder):
            os.remove(os.path.join(temp_image_folder, filename))
        os.rmdir(temp_image_folder)

        return StreamingResponse(pptx_io, media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation', headers={"Content-Disposition": "attachment; filename=compressed_presentation.pptx"})

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/pdf_to_images/")
async def pdf_to_images(file: UploadFile = File(...)):
    try:
        pdf_bytes = await file.read()
        pdf_document = fitz.open("pdf", pdf_bytes)

        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, 'w') as zipf:
            for i in range(len(pdf_document)):
                page = pdf_document[i]
                pix = page.get_pixmap()
                img_buffer = io.BytesIO()
                pix.save(img_buffer)
                img_buffer.seek(0)
                zipf.writestr(f"page_{i + 1}.png", img_buffer.getvalue())

        pdf_document.close()
        zip_buffer.seek(0)

        return StreamingResponse(zip_buffer, media_type='application/zip', headers={"Content-Disposition": "attachment; filename=output_images.zip"})

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/images_to_pdf/")
async def images_to_pdf(files: List[UploadFile] = File(...)):
    try:
        images = []
        for file in files:
            ext = os.path.splitext(file.filename)[1].lower()
            if ext not in ['.jpg', '.jpeg', '.png']:
                raise HTTPException(status_code=400, detail=f"The file {file.filename} has an unsupported file extension.")

            img = Image.open(await file.read())
            if img.mode in ("RGBA", "P"):
                img = img.convert("RGB")
            images.append(img)

        output_pdf = io.BytesIO()
        images[0].save(output_pdf, save_all=True, append_images=images[1:])
        output_pdf.seek(0)

        return StreamingResponse(output_pdf, media_type='application/pdf', headers={"Content-Disposition": "attachment; filename=img_to_pdf.pdf"})

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/compress_docx/")
async def compress_docx(file: UploadFile = File(...)):
    try:
        output_stream = io.BytesIO()
        doc = DocxDocument(await file.read())

        initial_size = len(await file.read()) / 1024
        await file.seek(0)

        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                image = rel.target_part.blob
                image_stream = io.BytesIO(image)

                with Image.open(image_stream) as img:
                    compressed_image_stream = io.BytesIO()
                    img.save(compressed_image_stream, format="JPEG", quality=60)
                    compressed_image = compressed_image_stream.getvalue()
                    rel.target_part._blob = compressed_image

        doc.save(output_stream)
        output_stream.seek(0)

        final_size = output_stream.tell() / 1024
        compression_ratio = (1 - (final_size / initial_size)) * 100

        response_message = {
            "message": "File compressed successfully.",
            "initial_size_kb": initial_size,
            "final_size_kb": final_size,
            "compression_ratio_percent": compression_ratio
        }

        return StreamingResponse(output_stream, media_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document', headers={"Content-Disposition": f'attachment; filename="compressed_{file.filename}"', "X-Message": str(response_message)})

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/compress_excel/")
async def compress_excel(file: UploadFile = File(...)):
    try:
        output_stream = io.BytesIO()
        workbook = openpyxl.load_workbook(await file.read())

        initial_size = len(await file.read()) / 1024
        await file.seek(0)

        for sheet in workbook.worksheets:
            for image in sheet._images:
                img_stream = io.BytesIO()
                image.ref.save(img_stream, format="PNG")
                img_stream.seek(0)

                with Image.open(img_stream) as img:
                    compressed_img_stream = io.BytesIO()
                    img.save(compressed_img_stream, format="JPEG", quality=80)
                    compressed_img_stream.seek(0)

                    compressed_img = OpenpyxlImage(compressed_img_stream)
                    compressed_img.anchor = image.anchor
                    sheet.add_image(compressed_img)
                    sheet._images.remove(image)

        workbook.save(output_stream)
        output_stream.seek(0)

        final_size = output_stream.tell() / 1024
        compression_ratio = (1 - (final_size / initial_size)) * 100

        response_message = {
            "message": "File compressed successfully.",
            "initial_size_kb": initial_size,
            "final_size_kb": final_size,
            "compression_ratio_percent": compression_ratio
        }

        return StreamingResponse(output_stream, media_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', headers={"Content-Disposition": f'attachment; filename="compressed_{file.filename}"', "X-Message": str(response_message)})

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/compress_pdf/")
async def compress_pdf(file: UploadFile = File(...)):
    try:
        output_stream = io.BytesIO()
        pdf_reader = PdfReader(await file.read())
        original_size = len(await file.read())
        await file.seek(0)

        pdf_writer = PdfWriter()

        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            page.compress_content_streams()
            pdf_writer.add_page(page)

        pdf_writer.write(output_stream)
        compressed_size = output_stream.tell()
        output_stream.seek(0)

        response_message = {
            "message": "File compressed successfully.",
            "original_size_bytes": original_size,
            "compressed_size_bytes": compressed_size,
            "size_reduced_bytes": original_size - compressed_size
        }

        return StreamingResponse(output_stream, media_type='application/pdf', headers={"Content-Disposition": f'attachment; filename="compressed_{file.filename}"', "X-Message": str(response_message)})

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    
if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=4050)



